"""
CHARAMOU AI - KnowledgeManager v2
Ingestion : PDF, DOCX, XLSX, PPTX, TXT, HTML, Markdown, Web
Pipeline RAG : chunk → index → recherche sémantique → réponse augmentée
"""
import os, re, json
from pathlib import Path
from typing import List, Dict, Optional, Any
from core.logger import setup_logger

logger = setup_logger("KnowledgeManager")
KNOWLEDGE_DIR = Path(__file__).parent.parent.parent / "data" / "knowledge"

# ── Parseurs par extension ────────────────────────────────────────────────────
SUPPORTED = {
    ".txt":  "text",
    ".md":   "markdown",
    ".html": "html",
    ".htm":  "html",
    ".pdf":  "pdf",
    ".docx": "docx",
    ".doc":  "docx",
    ".xlsx": "xlsx",
    ".xls":  "xlsx",
    ".pptx": "pptx",
    ".ppt":  "pptx",
    ".csv":  "csv",
    ".json": "json",
}


class KnowledgeManager:
    """
    Gère une base de connaissances locale.
    Supporte : PDF, DOCX, XLSX, PPTX, HTML, Markdown, TXT, CSV, JSON, Web.
    """

    CHUNK_SIZE    = 500
    CHUNK_OVERLAP = 100

    def __init__(self, memory=None, ai_client=None):
        self.memory    = memory
        self.ai_client = ai_client
        KNOWLEDGE_DIR.mkdir(parents=True, exist_ok=True)
        self._sources: List[Dict] = []
        logger.info("KnowledgeManager v2 initialisé.")

    # ── Entrée principale ─────────────────────────────────────────────────────

    def ingest_file(self, file_path: str) -> Dict:
        path   = Path(file_path)
        if not path.exists():
            return {"error": f"Fichier introuvable : {file_path}"}
        suffix = path.suffix.lower()
        fmt    = SUPPORTED.get(suffix)
        if not fmt:
            return {"error": f"Format non supporté '{suffix}'. Supportés : {list(SUPPORTED.keys())}"}

        logger.info(f"Ingestion : '{path.name}' ({fmt})")
        parser = getattr(self, f"_read_{fmt}", None)
        if not parser:
            return {"error": f"Parseur '{fmt}' non implémenté."}

        try:
            content = parser(str(path))
        except Exception as e:
            logger.error(f"Erreur lecture '{path.name}' : {e}")
            return {"error": str(e)}

        if not content or not content.strip():
            return {"error": "Document vide ou illisible."}

        return self._index_content(content, path.name, tags=[suffix, fmt])

    def ingest_text(self, text: str, source: str = "manuel", tags: list = None) -> Dict:
        return self._index_content(text, source, tags=tags or [])

    def ingest_url(self, url: str) -> Dict:
        try:
            import requests
            from html.parser import HTMLParser

            class Extractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self._parts, self._skip = [], False
                def handle_starttag(self, t, a):
                    if t in ('script','style','nav','footer','head'): self._skip = True
                def handle_endtag(self, t):
                    if t in ('script','style','nav','footer','head'): self._skip = False
                def handle_data(self, d):
                    if not self._skip and d.strip(): self._parts.append(d.strip())

            r = requests.get(url, headers={"User-Agent": "CHARAMOU-AI/2.0"}, timeout=10)
            p = Extractor(); p.feed(r.text)
            text = " ".join(p._parts)
            return self._index_content(text, url, tags=["web", "url"])
        except Exception as e:
            return {"error": str(e)}

    # ── Parseurs ──────────────────────────────────────────────────────────────

    def _read_text(self, path: str) -> str:
        return Path(path).read_text(encoding="utf-8", errors="ignore")

    def _read_markdown(self, path: str) -> str:
        text = Path(path).read_text(encoding="utf-8", errors="ignore")
        # Supprime les balises Markdown pour l'indexation
        text = re.sub(r'#{1,6}\s+', '', text)
        text = re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', text)
        text = re.sub(r'`{1,3}[^`]*`{1,3}', '', text)
        text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
        return text

    def _read_html(self, path: str) -> str:
        from html.parser import HTMLParser
        class P(HTMLParser):
            def __init__(self):
                super().__init__()
                self._p, self._s = [], False
            def handle_starttag(self, t, a):
                if t in ('script','style'): self._s = True
            def handle_endtag(self, t):
                if t in ('script','style'): self._s = False
            def handle_data(self, d):
                if not self._s and d.strip(): self._p.append(d.strip())
        parser = P()
        parser.feed(Path(path).read_text(encoding="utf-8", errors="ignore"))
        return " ".join(parser._p)

    def _read_pdf(self, path: str) -> str:
        # Essai 1 : pdfplumber (meilleure qualité)
        try:
            import pdfplumber
            parts = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t: parts.append(t)
            return "\n\n".join(parts)
        except ImportError:
            pass
        # Essai 2 : PyPDF2
        try:
            import PyPDF2
            parts = []
            with open(path, 'rb') as f:
                for page in PyPDF2.PdfReader(f).pages:
                    t = page.extract_text()
                    if t: parts.append(t)
            return "\n\n".join(parts)
        except ImportError:
            return "[Installez pdfplumber ou PyPDF2 pour lire les PDF]"

    def _read_docx(self, path: str) -> str:
        try:
            from docx import Document
            doc = Document(path)
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            # Tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(c.text for c in row.cells if c.text.strip())
                    if row_text: parts.append(row_text)
            return "\n".join(parts)
        except ImportError:
            return "[python-docx requis : pip install python-docx]"

    def _read_xlsx(self, path: str) -> str:
        try:
            import openpyxl
            wb   = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Feuille: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_str = " | ".join(str(c) for c in row if c is not None)
                    if row_str.strip(): parts.append(row_str)
            return "\n".join(parts)
        except ImportError:
            # Fallback csv-like avec xlrd
            try:
                import xlrd
                wb   = xlrd.open_workbook(path)
                parts = []
                for sheet in wb.sheets():
                    parts.append(f"[Feuille: {sheet.name}]")
                    for r in range(sheet.nrows):
                        row_str = " | ".join(str(sheet.cell(r, c).value)
                                             for c in range(sheet.ncols))
                        parts.append(row_str)
                return "\n".join(parts)
            except ImportError:
                return "[openpyxl requis : pip install openpyxl]"

    def _read_pptx(self, path: str) -> str:
        try:
            from pptx import Presentation
            prs   = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"[Diapo {i+1}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            return "\n".join(parts)
        except ImportError:
            return "[python-pptx requis : pip install python-pptx]"

    def _read_csv(self, path: str) -> str:
        import csv
        rows = []
        with open(path, newline='', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(row))
        return "\n".join(rows)

    def _read_json(self, path: str) -> str:
        try:
            data = json.loads(Path(path).read_text(encoding="utf-8"))
            return json.dumps(data, ensure_ascii=False, indent=2)
        except Exception as e:
            return f"JSON illisible : {e}"

    # ── Indexation & RAG ─────────────────────────────────────────────────────

    def _index_content(self, content: str, source: str, tags: list = None) -> Dict:
        chunks = self._chunk(content)
        if self.memory:
            for i, chunk in enumerate(chunks):
                self.memory.add_knowledge(
                    source=f"{source}#{i}", content=chunk,
                    summary=chunk[:80], tags=tags or []
                )
        doc = {"source": source, "chunks": len(chunks), "chars": len(content)}
        self._sources.append(doc)
        logger.info(f"'{source}' indexé : {len(chunks)} chunks, {len(content)} chars")
        return doc

    def _chunk(self, text: str) -> List[str]:
        chunks, start = [], 0
        text = text.strip()
        while start < len(text):
            chunk = text[start:start + self.CHUNK_SIZE]
            if chunk.strip():
                chunks.append(chunk)
            start += self.CHUNK_SIZE - self.CHUNK_OVERLAP
        return chunks

    # ── Requête RAG ───────────────────────────────────────────────────────────

    def query(self, question: str, top_k: int = 3) -> str:
        if not self.memory:
            return "Base de connaissances non disponible."
        results = self.memory.search_knowledge(question, top_k=top_k)
        if not results:
            return f"Aucune information trouvée pour : « {question} »"

        context = "\n\n---\n\n".join(
            f"[{r['source']}]\n{r['content']}" for r in results
        )

        if self.ai_client and self.ai_client.available:
            prompt = (
                f"En utilisant uniquement ces informations, réponds en français :\n\n"
                f"CONTEXTE :\n{context}\n\n"
                f"QUESTION : {question}\n\nRÉPONSE COURTE :"
            )
            try:
                return self.ai_client.ask(prompt)
            except Exception as e:
                logger.warning(f"IA RAG : {e}")

        best = results[0]
        return f"D'après {best['source']} : {best['content'][:300]}"

    def search(self, query: str, top_k: int = 5) -> List[Dict]:
        if not self.memory: return []
        return self.memory.search_knowledge(query, top_k=top_k)

    def get_sources(self) -> List[Dict]: return self._sources

    def handle(self, entities: dict = None, context=None) -> str:
        query = (entities or {}).get("raw_text", "")
        if not query:
            return "Que souhaitez-vous chercher dans la base de connaissances ?"
        return self.query(query)

    @property
    def supported_formats(self) -> List[str]:
        return list(SUPPORTED.keys())
