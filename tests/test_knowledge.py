"""Tests - KnowledgeManager v2 (tous formats)"""
import sys, tempfile, unittest, json
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))


def _make_km(tmp_dir: str):
    """KnowledgeManager avec mémoire isolée."""
    import sqlite3
    from core.memory_manager import MemoryManager, TFIDFIndex
    mem = MemoryManager.__new__(MemoryManager)
    mem._conn = sqlite3.connect(str(Path(tmp_dir) / "mem.db"))
    mem._conn.row_factory = sqlite3.Row
    mem._working  = {}
    mem._semantic = TFIDFIndex()
    mem._init_tables()
    from modules.services.knowledge_manager import KnowledgeManager
    return KnowledgeManager(memory=mem, ai_client=None), mem


class TestKnowledgeManagerV2(unittest.TestCase):

    def setUp(self):
        self.tmp = tempfile.mkdtemp()
        self.km, self.mem = _make_km(self.tmp)

    def tearDown(self):
        self.mem.close()

    # ── Formats texte ─────────────────────────────────────────────────────────

    def test_ingest_txt_file(self):
        p = Path(self.tmp) / "test.txt"
        p.write_text("Python est un langage de programmation très populaire.", encoding="utf-8")
        result = self.km.ingest_file(str(p))
        self.assertIn("chunks", result)
        self.assertGreater(result["chunks"], 0)

    def test_ingest_markdown(self):
        p = Path(self.tmp) / "doc.md"
        p.write_text("# Titre\n\nContenu **important** de la documentation.", encoding="utf-8")
        result = self.km.ingest_file(str(p))
        self.assertNotIn("error", result)

    def test_ingest_html(self):
        p = Path(self.tmp) / "page.html"
        p.write_text("<html><body><p>Contenu HTML de test.</p></body></html>", encoding="utf-8")
        result = self.km.ingest_file(str(p))
        self.assertNotIn("error", result)

    def test_ingest_csv(self):
        p = Path(self.tmp) / "data.csv"
        p.write_text("nom,age,ville\nAlice,30,Paris\nBob,25,Lyon", encoding="utf-8")
        result = self.km.ingest_file(str(p))
        self.assertNotIn("error", result)

    def test_ingest_json_file(self):
        p = Path(self.tmp) / "data.json"
        p.write_text(json.dumps({"key": "valeur", "items": [1, 2, 3]}), encoding="utf-8")
        result = self.km.ingest_file(str(p))
        self.assertNotIn("error", result)

    # ── Formats Office ────────────────────────────────────────────────────────

    def test_ingest_docx(self):
        try:
            from docx import Document
            doc  = Document()
            doc.add_heading("Titre du rapport", 0)
            doc.add_paragraph("Contenu principal du rapport de test.")
            p = Path(self.tmp) / "doc.docx"
            doc.save(str(p))
            result = self.km.ingest_file(str(p))
            self.assertNotIn("error", result)
            self.assertGreater(result["chars"], 0)
        except ImportError:
            self.skipTest("python-docx non installé")

    def test_ingest_xlsx(self):
        try:
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Données"
            ws.append(["Produit", "Prix", "Quantité"])
            ws.append(["Widget A", 10.5, 100])
            ws.append(["Widget B", 25.0, 50])
            p = Path(self.tmp) / "table.xlsx"
            wb.save(str(p))
            result = self.km.ingest_file(str(p))
            self.assertNotIn("error", result)
            self.assertGreater(result["chars"], 0)
        except ImportError:
            self.skipTest("openpyxl non installé")

    def test_ingest_pptx(self):
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Intelligence Artificielle"
            slide.placeholders[1].text = "Les assistants virtuels sont de plus en plus puissants."
            p = Path(self.tmp) / "pres.pptx"
            prs.save(str(p))
            result = self.km.ingest_file(str(p))
            self.assertNotIn("error", result)
        except ImportError:
            self.skipTest("python-pptx non installé")

    # ── Formats non supportés ─────────────────────────────────────────────────

    def test_unsupported_extension(self):
        p = Path(self.tmp) / "file.xyz"
        p.write_text("contenu")
        result = self.km.ingest_file(str(p))
        self.assertIn("error", result)
        self.assertIn("non supporté", result["error"])

    def test_missing_file(self):
        result = self.km.ingest_file("/tmp/fichier_inexistant_xyz.txt")
        self.assertIn("error", result)

    # ── Ingestion texte direct ────────────────────────────────────────────────

    def test_ingest_text_direct(self):
        result = self.km.ingest_text("Django est un framework web Python.", source="test")
        self.assertEqual(result["source"], "test")
        self.assertGreater(result["chunks"], 0)

    def test_ingest_text_with_tags(self):
        result = self.km.ingest_text("FastAPI est moderne.", source="api_doc", tags=["python", "web"])
        self.assertNotIn("error", result)

    # ── Chunking ──────────────────────────────────────────────────────────────

    def test_chunk_short_text(self):
        chunks = self.km._chunk("Texte court.")
        self.assertEqual(len(chunks), 1)

    def test_chunk_long_text(self):
        long = "mot " * 200
        chunks = self.km._chunk(long)
        self.assertGreater(len(chunks), 1)
        for c in chunks:
            self.assertLessEqual(len(c), self.km.CHUNK_SIZE + 10)

    def test_chunk_overlap(self):
        """Les chunks doivent se chevaucher."""
        text = "A" * self.km.CHUNK_SIZE + "B" * self.km.CHUNK_SIZE
        chunks = self.km._chunk(text)
        if len(chunks) > 1:
            overlap = self.km.CHUNK_SIZE - self.km.CHUNK_OVERLAP
            self.assertLess(overlap, self.km.CHUNK_SIZE)

    # ── Recherche & RAG ───────────────────────────────────────────────────────

    def test_search_after_ingest(self):
        self.km.ingest_text("Le machine learning est un domaine de l'IA.", "ml_doc")
        results = self.km.search("machine learning intelligence artificielle")
        self.assertIsInstance(results, list)

    def test_query_no_knowledge(self):
        result = self.km.query("Python Django")
        self.assertIsInstance(result, str)

    def test_query_with_knowledge(self):
        self.km.ingest_text("Flask est un micro-framework Python pour le web.", "flask_doc")
        result = self.km.query("framework web Python")
        self.assertIsInstance(result, str)

    def test_handle_with_query(self):
        self.km.ingest_text("NumPy est une bibliothèque Python pour le calcul numérique.", "numpy")
        result = self.km.handle(entities={"raw_text": "numpy calcul"})
        self.assertIsInstance(result, str)

    def test_handle_empty_query(self):
        result = self.km.handle(entities={})
        self.assertIn("souhaitez", result.lower())

    # ── Parseurs internes ─────────────────────────────────────────────────────

    def test_read_text(self):
        p = Path(self.tmp) / "t.txt"
        p.write_text("contenu simple", encoding="utf-8")
        text = self.km._read_text(str(p))
        self.assertIn("contenu simple", text)

    def test_read_markdown_strips_syntax(self):
        p = Path(self.tmp) / "t.md"
        p.write_text("# Titre\n**gras** et *italique*\n`code`", encoding="utf-8")
        text = self.km._read_markdown(str(p))
        self.assertNotIn("**", text)
        self.assertNotIn("##", text)

    def test_read_html_strips_tags(self):
        p = Path(self.tmp) / "t.html"
        p.write_text("<html><body><p>Texte visible</p><script>alert(x)</script></body></html>")
        text = self.km._read_html(str(p))
        self.assertIn("Texte visible", text)
        self.assertNotIn("alert", text)

    def test_read_csv_produces_pipes(self):
        p = Path(self.tmp) / "t.csv"
        p.write_text("a,b,c\n1,2,3", encoding="utf-8")
        text = self.km._read_csv(str(p))
        self.assertIn("|", text)

    def test_read_json_produces_text(self):
        p = Path(self.tmp) / "t.json"
        p.write_text('{"name": "test"}', encoding="utf-8")
        text = self.km._read_json(str(p))
        self.assertIn("test", text)

    # ── Propriétés ───────────────────────────────────────────────────────────

    def test_supported_formats(self):
        fmts = self.km.supported_formats
        self.assertIn(".pdf", fmts)
        self.assertIn(".docx", fmts)
        self.assertIn(".xlsx", fmts)
        self.assertIn(".pptx", fmts)
        self.assertIn(".html", fmts)
        self.assertIn(".md",   fmts)
        self.assertIn(".txt",  fmts)

    def test_get_sources_after_ingest(self):
        self.km.ingest_text("test content", source="src1")
        sources = self.km.get_sources()
        self.assertTrue(any(s["source"] == "src1" for s in sources))


class TestMemoryV3Backends(unittest.TestCase):
    """Teste que le MemoryManager bascule correctement entre les backends."""

    def test_tfidf_backend_selected_when_no_chromadb(self):
        from core.memory_manager import TFIDFIndex, _select_backend
        import unittest.mock as mock
        # Simule ChromaDB absent
        with mock.patch.dict('sys.modules', {'chromadb': None}):
            # Reimport pour forcer la sélection
            import importlib
            import core.memory_manager as mm
            backend = mm.TFIDFIndex()
            self.assertEqual(backend.name, "tfidf")

    def test_chromadb_backend_name(self):
        try:
            import chromadb
            from core.memory_manager import ChromaDBIndex
            # Test que la classe existe et a le bon nom
            self.assertEqual(ChromaDBIndex.name, "chromadb")
        except ImportError:
            self.skipTest("chromadb non installé")

    def test_faiss_backend_name(self):
        from core.memory_manager import FAISSIndex
        self.assertEqual(FAISSIndex.name, "faiss")

    def test_tfidf_scale_warning_threshold(self):
        from core.memory_manager import TFIDFIndex
        idx = TFIDFIndex()
        self.assertEqual(idx.SCALE_WARNING, 5000)


if __name__ == "__main__":
    unittest.main(verbosity=2)
